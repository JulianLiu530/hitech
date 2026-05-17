#!/usr/bin/env python3
"""Generate HiTech_Group2_Sharing.pptx from extracted slide content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ────────────────────────────────────────────────────────────
BLUE       = RGBColor(0x00, 0x71, 0xE3)   # Apple blue
DARK       = RGBColor(0x1D, 0x1D, 0x1F)   # Near-black
GRAY       = RGBColor(0x6E, 0x6E, 0x73)   # Mid-gray
GRAY_LIGHT = RGBColor(0x86, 0x86, 0x8B)
BG_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BG_GRAY    = RGBColor(0xF5, 0xF5, 0xF7)
GREEN      = RGBColor(0x1A, 0x87, 0x33)
RED        = RGBColor(0xC0, 0x39, 0x2B)
ORANGE     = RGBColor(0xB3, 0x62, 0x00)

# ── Slide size: 16:9 widescreen ───────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── Helper functions ──────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()  # no line by default
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.5)
    return shape

def add_text(slide, text, l, t, w, h,
             size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def add_para(tf, text, size=12, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, space_before=0, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p

def eyebrow(slide, text, t_offset=Inches(0.55)):
    add_text(slide, text,
             l=Inches(1.1), t=t_offset, w=Inches(11), h=Inches(0.35),
             size=9, bold=True, color=BLUE)

def slide_bg(slide, color):
    bg = add_rect(slide, 0, 0, W, H, fill=color)
    # push to back
    sp = bg._element
    slide.shapes._spTree.remove(sp)
    slide.shapes._spTree.insert(2, sp)

def h_line(slide, t):
    add_rect(slide, Inches(1.1), t, Inches(11.13), Pt(1),
             fill=RGBColor(0xD2, 0xD2, 0xD7))

def pill_box(slide, text, l, t, color_fill, color_text):
    box = add_rect(slide, l, t, Inches(1.5), Inches(0.3), fill=color_fill)
    add_text(slide, text,
             l=l, t=t, w=Inches(1.5), h=Inches(0.3),
             size=8, bold=True, color=color_text, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 1 — COVER
# ─────────────────────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
slide_bg(s1, BG_WHITE)

# Left accent strip
add_rect(s1, 0, 0, Inches(0.06), H, fill=BLUE)

# Eyebrow
add_text(s1, "MGNT 5506GA  ·  Strategic Management  ·  CUHK",
         Inches(1.1), Inches(0.6), Inches(8), Inches(0.4),
         size=9, bold=True, color=BLUE)

# Main title
add_text(s1, "Digital",
         Inches(1.1), Inches(1.3), Inches(7), Inches(1.1),
         size=72, bold=True, color=DARK)
add_text(s1, "Blind Spots",
         Inches(1.1), Inches(2.25), Inches(7), Inches(1.1),
         size=72, bold=True, color=DARK)

# Subtitle
add_text(s1, "A Strategic Analysis of HiTech Techno Limited —\nNavigating the Digital Transformation Imperative",
         Inches(1.1), Inches(3.5), Inches(7), Inches(0.9),
         size=14, color=GRAY)

# Meta row
add_text(s1, "Group:  Group 2 — Circuit Minds     Date: April 17, 2026     Duration: 8 Minutes",
         Inches(1.1), Inches(4.55), Inches(10), Inches(0.4),
         size=11, color=DARK)

# Divider
h_line(s1, Inches(5.1))

# Team members
add_text(s1, "TEAM   Daria Iudicheva  ·  Zhi Niu  ·  Qiaochu Huang  ·  Bhavya Jain  ·  Junqi Pu  ·  Zhilang Liu",
         Inches(1.1), Inches(5.25), Inches(11), Inches(0.4),
         size=10, color=GRAY)

# Company card (right side)
add_rect(s1, Inches(9.5), Inches(1.3), Inches(2.8), Inches(3.5),
         fill=BG_WHITE, line=RGBColor(0xD2, 0xD2, 0xD7))
add_text(s1, "HiTech Techno\nLimited",
         Inches(9.55), Inches(1.45), Inches(2.7), Inches(1.2),
         size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_text(s1, "B2B IT Component Distributor\nEst. 2018  ·  TST, Hong Kong",
         Inches(9.55), Inches(3.0), Inches(2.7), Inches(0.7),
         size=9, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 2 — BUSINESS INTRODUCTION
# ─────────────────────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
slide_bg(s2, BG_GRAY)
add_rect(s2, 0, 0, Inches(0.06), H, fill=BLUE)

eyebrow(s2, "01 — Business Introduction")
add_text(s2, "Getting Inside HiTech Techno",
         Inches(1.1), Inches(0.9), Inches(11), Inches(0.75),
         size=34, bold=True, color=DARK)
h_line(s2, Inches(1.72))

# ── Left column: Company Profile ─────────────────────────────────
add_rect(s2, Inches(1.1), Inches(1.85), Inches(5.5), Inches(2.5),
         fill=BG_WHITE, line=RGBColor(0xD2, 0xD2, 0xD7))
add_text(s2, "COMPANY PROFILE",
         Inches(1.2), Inches(1.95), Inches(5.3), Inches(0.3),
         size=8, bold=True, color=BLUE)

profile_lines = [
    "B2B IT component distributor, est. 2018",
    "Headquartered in Tsim Sha Tsui, Kowloon, HK",
    "Sources from tier-1 global manufacturers",
    "Serves corporate IT buyers, system integrators & resellers",
    "~40 active accounts via personal referral networks",
]
for i, line in enumerate(profile_lines):
    add_text(s2, f"• {line}",
             Inches(1.2), Inches(2.3) + Inches(0.33)*i, Inches(5.2), Inches(0.35),
             size=10, color=DARK)

# ── Left column: Product Portfolio ───────────────────────────────
add_rect(s2, Inches(1.1), Inches(4.5), Inches(5.5), Inches(2.6),
         fill=BG_WHITE, line=RGBColor(0xD2, 0xD2, 0xD7))
add_text(s2, "PRODUCT PORTFOLIO",
         Inches(1.2), Inches(4.6), Inches(5.3), Inches(0.3),
         size=8, bold=True, color=BLUE)

products = ["SSDs & NVMe Drives", "RAM Modules (DDR3/4/5)",
            "Flash Memory", "CPUs & Motherboards",
            "Storage Devices", "GPUs (Gaming / AI / HPC)"]
for i, p in enumerate(products):
    col = 0 if i < 3 else 1
    row = i % 3
    add_text(s2, f"• {p}",
             Inches(1.2) + Inches(2.7)*col, Inches(5.0) + Inches(0.33)*row,
             Inches(2.6), Inches(0.35), size=10, color=DARK)

# ── Right column: How we got access ──────────────────────────────
add_rect(s2, Inches(6.9), Inches(1.85), Inches(5.2), Inches(2.0),
         fill=BG_WHITE, line=RGBColor(0xD2, 0xD2, 0xD7))
add_text(s2, "HOW WE GOT ACCESS",
         Inches(7.0), Inches(1.95), Inches(5.0), Inches(0.3),
         size=8, bold=True, color=BLUE)
add_text(s2, "Bhavya Jain (group member) has direct family ties to\nHiTech Techno's ownership — a family business.\n\nThis enabled privileged access to management, operations\nstaff, and internal documents, allowing us to observe actual\nworkflows rather than relying on reported practices.",
         Inches(7.0), Inches(2.3), Inches(5.0), Inches(1.4),
         size=10, color=DARK)

# ── Right column: Fieldwork stats ───────────────────────────────
stats = [("2", "On-site Visits"), ("4", "Staff Interviews"),
         ("3", "Docs Reviewed"), ("5", "Competitors")]
for i, (num, label) in enumerate(stats):
    x = Inches(6.9) + Inches(1.3)*i
    add_rect(s2, x, Inches(4.05), Inches(1.2), Inches(0.9),
             fill=BG_WHITE, line=RGBColor(0xD2, 0xD2, 0xD7))
    add_text(s2, num, x, Inches(4.08), Inches(1.2), Inches(0.45),
             size=26, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(s2, label, x, Inches(4.55), Inches(1.2), Inches(0.35),
             size=8, color=GRAY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 3 — FIELD OBSERVATIONS
# ─────────────────────────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
slide_bg(s3, BG_WHITE)
add_rect(s3, 0, 0, Inches(0.06), H, fill=BLUE)

eyebrow(s3, "02 — Field Observations")
add_text(s3, "What We Saw On the Ground",
         Inches(1.1), Inches(0.9), Inches(11), Inches(0.75),
         size=34, bold=True, color=DARK)
h_line(s3, Inches(1.72))

# Visit 1 card
add_rect(s3, Inches(1.1), Inches(1.85), Inches(5.5), Inches(2.8),
         fill=BG_GRAY, line=RGBColor(0xD2, 0xD2, 0xD7))
add_text(s3, "VISIT 1 — March 18, 2026",
         Inches(1.2), Inches(1.95), Inches(5.3), Inches(0.3),
         size=8, bold=True, color=BLUE)
v1_items = [
    "Full office walkthrough & workflow observation",
    "Orders managed via Microsoft Excel — no CRM or ERP",
    "Client comms via WhatsApp & email only",
    "Interviewed Mr. Chan (Ops Manager) & Sales Lead",
    "No digital marketing assets or brand presence",
]
for i, item in enumerate(v1_items):
    add_text(s3, f"• {item}",
             Inches(1.2), Inches(2.35) + Inches(0.42)*i, Inches(5.2), Inches(0.4),
             size=10, color=DARK)

# Visit 2 card
add_rect(s3, Inches(6.9), Inches(1.85), Inches(5.2), Inches(2.8),
         fill=BG_GRAY, line=RGBColor(0xD2, 0xD2, 0xD7))
add_text(s3, "VISIT 2 — April 2, 2026",
         Inches(7.0), Inches(1.95), Inches(5.0), Inches(0.3),
         size=8, bold=True, color=BLUE)
v2_items = [
    "Observed a live client negotiation call",
    "Catalog distributed as static PDFs — no live pricing",
    "Zero live inventory system — tracked on paper & Excel",
    "Interviewed Procurement Officer & Admin Assistant",
    "Reviewed 3 internal docs: sales log, supplier list, templates",
]
for i, item in enumerate(v2_items):
    add_text(s3, f"• {item}",
             Inches(7.0), Inches(2.35) + Inches(0.42)*i, Inches(5.0), Inches(0.4),
             size=10, color=DARK)

# Key stats strip
kstats = [("120–150", "Orders/Month"), ("80%", "Repeat Revenue"),
          ("0", "Social Media"), ("<200", "Website Visits/Mo")]
for i, (num, label) in enumerate(kstats):
    x = Inches(1.1) + Inches(2.83)*i
    add_rect(s3, x, Inches(4.85), Inches(2.7), Inches(1.0),
             fill=BG_GRAY, line=RGBColor(0xD2, 0xD2, 0xD7))
    add_text(s3, num, x, Inches(4.9), Inches(2.7), Inches(0.55),
             size=28, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(s3, label, x, Inches(5.5), Inches(2.7), Inches(0.3),
             size=9, color=GRAY, align=PP_ALIGN.CENTER)

# Timeline
timeline = [
    ("Mar 18", "Site Visit 1"),
    ("Mar 25", "Analysis & Benchmarking"),
    ("Apr 2",  "Site Visit 2"),
    ("Apr 10", "Digital Audit"),
    ("Apr 17", "Group Sharing"),
]
add_text(s3, "FIELDWORK TIMELINE",
         Inches(1.1), Inches(6.05), Inches(11), Inches(0.25),
         size=8, bold=True, color=BLUE)
for i, (date, event) in enumerate(timeline):
    x = Inches(1.1) + Inches(2.4)*i
    add_text(s3, date, x, Inches(6.35), Inches(2.3), Inches(0.25),
             size=9, bold=True, color=BLUE)
    add_text(s3, event, x, Inches(6.62), Inches(2.3), Inches(0.3),
             size=9, color=DARK)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 4 — MAIN PROBLEM
# ─────────────────────────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
slide_bg(s4, BG_GRAY)
add_rect(s4, 0, 0, Inches(0.06), H, fill=BLUE)

# Eyebrow
eyebrow(s4, "03 — CORE PROBLEM IDENTIFIED")

# Two-tone title: red first line, dark second line (like reference image)
add_text(s4, "Digital Invisibility",
         Inches(1.1), Inches(0.78), Inches(11), Inches(0.65),
         size=36, bold=True, color=RED)
add_text(s4, "in a Digitally-Driven B2B Market",
         Inches(1.1), Inches(1.38), Inches(11), Inches(0.55),
         size=32, bold=True, color=DARK)

# ── Two main cards side by side ───────────────────────────────────────────────
CARD_L = Inches(1.1)
CARD_T = Inches(2.1)
CARD_W = Inches(5.3)
CARD_H = Inches(5.0)
TABLE_L = Inches(6.8)
TABLE_W = Inches(6.15)

# Left card background
add_rect(s4, CARD_L, CARD_T, CARD_W, CARD_H,
         fill=BG_WHITE, line=RGBColor(0xD2, 0xD2, 0xD7))

# ⚠ KEY DEFICIENCIES label
add_text(s4, "△  KEY DEFICIENCIES",
         CARD_L + Inches(0.2), CARD_T + Inches(0.15),
         Inches(4.8), Inches(0.3),
         size=8, bold=True, color=RED)

# Bullet items — bold keyword + plain detail on same/next line
deficiencies = [
    ("No social media presence",        "whatsoever — LinkedIn,\n        Facebook, Instagram: all absent"),
    ("Static brochure-only website",    "— no pricing, no search,\n        no live inventory, no contact form"),
    ("100% word-of-mouth acquisition",  "— zero inbound\n        digital leads ever recorded"),
    ("Zero digital marketing spend",    "— no Google Ads, no SEO,\n        no content strategy"),
    ("Invisible to new B2B buyers",     "— <200 website\n        visits/month vs. thousands for competitors"),
]

DOT_COLOR = RGBColor(0xFF, 0x3B, 0x30)   # red dot
for i, (bold_part, plain_part) in enumerate(deficiencies):
    item_y = CARD_T + Inches(0.58) + Inches(0.82)*i

    # Red bullet dot
    dot = s4.shapes.add_shape(1,
        CARD_L + Inches(0.15), item_y + Inches(0.09),
        Inches(0.12), Inches(0.12))
    dot.fill.solid(); dot.fill.fore_color.rgb = DOT_COLOR
    dot.line.fill.background()

    # Text box with bold+plain in one textbox
    txb = s4.shapes.add_textbox(
        CARD_L + Inches(0.35), item_y,
        Inches(4.7), Inches(0.75))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    # bold run
    r1 = p.add_run(); r1.text = bold_part
    r1.font.bold = True; r1.font.size = Pt(10.5)
    r1.font.color.rgb = DARK
    # plain run
    r2 = p.add_run(); r2.text = " " + plain_part.replace("\n        ", " ")
    r2.font.bold = False; r2.font.size = Pt(10.5)
    r2.font.color.rgb = DARK

# ── Right card: comparison table ─────────────────────────────────────────────
add_rect(s4, TABLE_L, CARD_T, TABLE_W, CARD_H,
         fill=BG_WHITE, line=RGBColor(0xD2, 0xD2, 0xD7))

# Column headers
col_labels   = ["DIMENSION", "HITECH", "TECHDATA\nHK", "INGRAM\nMICRO", "ARROW\nELEC."]
col_w        = [Inches(1.45), Inches(0.95), Inches(0.93), Inches(0.93), Inches(0.93)]
col_x        = [TABLE_L + Inches(0.06)]
for cw in col_w[:-1]:
    col_x.append(col_x[-1] + cw)

hdr_y  = CARD_T + Inches(0.1)
hdr_h  = Inches(0.45)

for j, (lbl, cx, cw) in enumerate(zip(col_labels, col_x, col_w)):
    c = RED if j == 1 else GRAY_LIGHT
    add_text(s4, lbl, cx, hdr_y, cw, hdr_h,
             size=7.5, bold=True, color=c, align=PP_ALIGN.CENTER)

# Row data — (dimension, hitech_val, t1, t2, t3)
# pill colours: RED_PILL for ✕, GREEN_PILL for ✓, AMBER_PILL for ●
RED_PILL    = RGBColor(0xFF, 0xE5, 0xE5)
GREEN_PILL  = RGBColor(0xE2, 0xF5, 0xE8)
AMBER_PILL  = RGBColor(0xFF, 0xF3, 0xD6)

table_rows = [
    ("LinkedIn\nPresence",  ("✕","None"),  ("✓","Active"), ("✓","Active"), ("✓","Active")),
    ("E-Commerce\nPortal",  ("✕","None"),  ("✓","Full"),   ("✓","Full"),   ("●","Partial")),
    ("Live Inventory",      ("✕","None"),  ("✓","Yes"),    ("✓","Yes"),    ("✓","Yes")),
    ("Digital\nMarketing",  ("✕","Zero"),  ("✓","Active"), ("✓","Active"), ("✓","Active")),
    ("Discoverability",     ("✕","<5%"),   ("✓","High"),   ("✓","Very\nHigh"), ("✓","High")),
]

ROW_H   = Inches(0.84)
row_y0  = CARD_T + hdr_h + Inches(0.15)

# Thin separator line under header
add_rect(s4, TABLE_L + Inches(0.06), row_y0 - Inches(0.04),
         TABLE_W - Inches(0.12), Pt(0.5),
         fill=RGBColor(0xD2, 0xD2, 0xD7))

for ri, (dim, *cells) in enumerate(table_rows):
    ry = row_y0 + ROW_H * ri
    # Thin row separator (except first)
    if ri > 0:
        add_rect(s4, TABLE_L + Inches(0.06), ry - Inches(0.04),
                 TABLE_W - Inches(0.12), Pt(0.5),
                 fill=RGBColor(0xE8, 0xE8, 0xED))

    # Dimension label
    add_text(s4, dim, col_x[0], ry + Inches(0.1),
             col_w[0], ROW_H - Inches(0.1),
             size=10, color=DARK, align=PP_ALIGN.LEFT)

    # Pill cells
    for ci, (icon, val) in enumerate(cells):
        cx   = col_x[ci + 1]
        cw   = col_w[ci + 1]
        pill_fill  = RED_PILL  if icon == "✕" else (AMBER_PILL if icon == "●" else GREEN_PILL)
        pill_color = RED       if icon == "✕" else (ORANGE     if icon == "●" else GREEN)

        pill_w = Inches(0.75)
        pill_h = Inches(0.6)
        pill_x = cx + (cw - pill_w) / 2
        pill_y = ry + (ROW_H - pill_h) / 2

        # Rounded pill box
        add_rect(s4, pill_x, pill_y, pill_w, pill_h, fill=pill_fill)

        # Icon + value inside pill
        add_text(s4, f"{icon}\n{val}",
                 pill_x, pill_y,
                 pill_w, pill_h,
                 size=9, bold=True, color=pill_color,
                 align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 5 — WHY STRATEGIC
# ─────────────────────────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
slide_bg(s5, BG_WHITE)
add_rect(s5, 0, 0, Inches(0.06), H, fill=BLUE)

eyebrow(s5, "04 — Strategic Framing")
add_text(s5, "Why This Is Not Just a Marketing Problem",
         Inches(1.1), Inches(0.9), Inches(11), Inches(0.75),
         size=30, bold=True, color=DARK)
h_line(s5, Inches(1.72))

# Porter's 5 Forces
add_text(s5, "PORTER'S FIVE FORCES",
         Inches(1.1), Inches(1.85), Inches(5.5), Inches(0.3),
         size=8, bold=True, color=BLUE)

forces = [
    ("Competitive Rivalry",   "HIGH RISK ▲▲",    RED),
    ("New Entrants",          "Medium ↑",         ORANGE),
    ("Threat of Substitutes", "HIGH ▲▲",          RED),
    ("Supplier Power",        "Medium",           ORANGE),
    ("Buyer Power",           "Medium-High",      ORANGE),
]
for i, (force, level, col) in enumerate(forces):
    y = Inches(2.2) + Inches(0.6)*i
    add_text(s5, force, Inches(1.1), y, Inches(3.0), Inches(0.55),
             size=10, color=DARK)
    add_text(s5, level, Inches(4.2), y, Inches(2.4), Inches(0.55),
             size=10, bold=True, color=col)

# Gartner quote
add_rect(s5, Inches(1.1), Inches(5.45), Inches(5.5), Inches(0.85),
         fill=RGBColor(0xE8, 0xF2, 0xFF), line=BLUE)
add_text(s5, '"80% of B2B sales interactions will occur through digital\nchannels by 2025"  — Gartner, 2020 Future of Sales',
         Inches(1.2), Inches(5.5), Inches(5.3), Inches(0.75),
         size=10, italic=True, color=DARK)

# RBV
add_text(s5, "RESOURCE-BASED VIEW (RBV)",
         Inches(7.0), Inches(1.85), Inches(5.2), Inches(0.3),
         size=8, bold=True, color=BLUE)

strengths = [
    ("✓ Relationships", "Deep trust-based network; 40+ clients; 80% repeat revenue"),
    ("✓ Sourcing",      "Established tier-1 manufacturer relationships; competitive pricing"),
]
gaps = [
    ("⚠ Digital Assets",  "Advantages invisible online. New buyers cannot discover HiTech."),
    ("⚠ Scalability",     "Non-digital resources cap growth at ~40 accounts."),
]

add_text(s5, "Strengths",
         Inches(7.0), Inches(2.25), Inches(5.2), Inches(0.3),
         size=9, bold=True, color=GREEN)
for i, (label, detail) in enumerate(strengths):
    y = Inches(2.6) + Inches(0.65)*i
    add_rect(s5, Inches(7.0), y, Inches(5.2), Inches(0.58),
             fill=RGBColor(0xEA, 0xF7, 0xEE), line=RGBColor(0xD2, 0xD2, 0xD7))
    add_text(s5, label, Inches(7.1), y + Inches(0.03), Inches(5.0), Inches(0.25),
             size=10, bold=True, color=GREEN)
    add_text(s5, detail, Inches(7.1), y + Inches(0.28), Inches(5.0), Inches(0.25),
             size=9, color=DARK)

add_text(s5, "Gaps",
         Inches(7.0), Inches(3.96), Inches(5.2), Inches(0.3),
         size=9, bold=True, color=RED)
for i, (label, detail) in enumerate(gaps):
    y = Inches(4.28) + Inches(0.65)*i
    add_rect(s5, Inches(7.0), y, Inches(5.2), Inches(0.58),
             fill=RGBColor(0xFF, 0xEE, 0xEC), line=RGBColor(0xD2, 0xD2, 0xD7))
    add_text(s5, label, Inches(7.1), y + Inches(0.03), Inches(5.0), Inches(0.25),
             size=10, bold=True, color=RED)
    add_text(s5, detail, Inches(7.1), y + Inches(0.28), Inches(5.0), Inches(0.25),
             size=9, color=DARK)

# Strategic risks
risks = ["Market share erosion to digital-native rivals",
         "Revenue plateau — capped at ~40 accounts",
         "Relationship attrition with no digital backup"]
add_text(s5, "STRATEGIC RISKS",
         Inches(7.0), Inches(5.55), Inches(5.2), Inches(0.3),
         size=8, bold=True, color=RED)
for i, r in enumerate(risks):
    add_text(s5, f"• {r}",
             Inches(7.0), Inches(5.9) + Inches(0.3)*i, Inches(5.2), Inches(0.28),
             size=10, color=DARK)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 6 — PROGRESS DASHBOARD
# ─────────────────────────────────────────────────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
slide_bg(s6, BG_GRAY)
add_rect(s6, 0, 0, Inches(0.06), H, fill=BLUE)

eyebrow(s6, "05 — Research Progress")
add_text(s6, "What We Have Done So Far",
         Inches(1.1), Inches(0.9), Inches(11), Inches(0.75),
         size=34, bold=True, color=DARK)
h_line(s6, Inches(1.72))

# Summary pills
for i, (n, label, col) in enumerate([
    ("5/8", "Tasks Done", GREEN),
    ("2",   "In Progress", ORANGE),
    ("1",   "Pending",    GRAY),
]):
    x = Inches(1.1) + Inches(2.0)*i
    add_rect(s6, x, Inches(1.82), Inches(1.8), Inches(0.75),
             fill=BG_WHITE, line=RGBColor(0xD2, 0xD2, 0xD7))
    add_text(s6, n, x, Inches(1.85), Inches(1.8), Inches(0.42),
             size=24, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s6, label, x, Inches(2.3), Inches(1.8), Inches(0.25),
             size=8, color=GRAY, align=PP_ALIGN.CENTER)

# Task table
tasks = [
    ("Site Visits (2/2)",                   100, GREEN),
    ("Staff Interviews (4/4)",              100, GREEN),
    ("Competitor Benchmarking (5 firms)",   100, GREEN),
    ("Website Traffic Analysis",            100, GREEN),
    ("Digital Audit (SEO & Social)",        100, GREEN),
    ("Financial Impact Modeling",            40, ORANGE),
    ("Client Survey (10 clients)",           30, ORANGE),
    ("Final Report & Recommendations",        0, GRAY),
]

add_text(s6, "TASK",
         Inches(1.1), Inches(2.75), Inches(5.5), Inches(0.3),
         size=8, bold=True, color=BLUE)
add_text(s6, "PROGRESS",
         Inches(6.8), Inches(2.75), Inches(2.0), Inches(0.3),
         size=8, bold=True, color=BLUE)
add_text(s6, "STATUS",
         Inches(9.0), Inches(2.75), Inches(1.5), Inches(0.3),
         size=8, bold=True, color=BLUE)

for i, (task, pct, col) in enumerate(tasks):
    y = Inches(3.1) + Inches(0.5)*i
    bg = BG_WHITE if i % 2 == 0 else BG_GRAY
    add_rect(s6, Inches(1.1), y, Inches(10.5), Inches(0.45),
             fill=bg, line=RGBColor(0xD2, 0xD2, 0xD7))
    add_text(s6, task, Inches(1.2), y + Inches(0.07), Inches(5.4), Inches(0.35),
             size=10, color=DARK)
    # progress bar track
    bar_x = Inches(6.8)
    bar_w = Inches(2.0)
    bar_h = Inches(0.12)
    bar_y = y + Inches(0.165)
    add_rect(s6, bar_x, bar_y, bar_w, bar_h,
             fill=RGBColor(0xE8, 0xE8, 0xED))
    if pct > 0:
        add_rect(s6, bar_x, bar_y, bar_w * pct / 100, bar_h, fill=col)
    add_text(s6, f"{pct}%", Inches(9.0), y + Inches(0.05), Inches(1.5), Inches(0.35),
             size=9, bold=True, color=col, align=PP_ALIGN.CENTER)

# Next steps
add_text(s6, "NEXT STEPS",
         Inches(10.8), Inches(2.75), Inches(2.3), Inches(0.3),
         size=8, bold=True, color=BLUE)
nextsteps = [
    ("Apr 17", "Group sharing — today"),
    ("Late Apr", "Financial model & client survey"),
    ("May",     "Final strategic recommendations"),
    ("Final",   "Submit to Prof. Sissi Li"),
]
for i, (date, step) in enumerate(nextsteps):
    y = Inches(3.1) + Inches(0.5)*i
    add_text(s6, date, Inches(10.8), y + Inches(0.05), Inches(0.85), Inches(0.35),
             size=8, bold=True, color=BLUE)
    add_text(s6, step, Inches(11.7), y + Inches(0.05), Inches(1.5), Inches(0.35),
             size=8, color=DARK)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 7 — THANK YOU
# ─────────────────────────────────────────────────────────────────────────────
s7 = prs.slides.add_slide(BLANK)
slide_bg(s7, BG_WHITE)
add_rect(s7, 0, 0, Inches(0.06), H, fill=BLUE)
# Bottom accent
add_rect(s7, 0, H - Inches(0.06), W, Inches(0.06), fill=BLUE)

add_text(s7, "Thank You",
         Inches(1.5), Inches(1.8), Inches(10), Inches(1.6),
         size=72, bold=True, color=DARK, align=PP_ALIGN.CENTER)

add_text(s7, "We welcome your questions & feedback",
         Inches(1.5), Inches(3.5), Inches(10), Inches(0.55),
         size=18, color=GRAY, align=PP_ALIGN.CENTER)

h_line(s7, Inches(4.3))

# Contact info
add_text(s7, "Circuit Minds  ·  Group 2  ·  MGNT 5506GA  ·  CUHK",
         Inches(1.5), Inches(4.55), Inches(10), Inches(0.4),
         size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

add_text(s7, "hitechtechnoltd@gmail.com   |   TST, Kowloon, Hong Kong   |   April 17, 2026",
         Inches(1.5), Inches(5.05), Inches(10), Inches(0.4),
         size=10, color=GRAY, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/llll/Desktop/HiTech_Group2_Sharing.pptx"
prs.save(out)
print(f"✅  Saved → {out}")
